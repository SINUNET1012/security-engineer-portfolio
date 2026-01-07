from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REPO_ROOT = ROOT.parents[1]


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


@dataclass(frozen=True)
class Theme:
    bg: RGBColor = rgb("#FFFFFF")
    text: RGBColor = rgb("#111827")
    muted: RGBColor = rgb("#374151")
    line: RGBColor = rgb("#E5E7EB")
    accent: RGBColor = rgb("#6E3CBC")  # purple
    accent_dark: RGBColor = rgb("#4C1D95")

    font: str = "Malgun Gothic"


THEME = Theme()


def render_pdf_page(
    pdf_path: Path,
    page_no_1_based: int,
    out_path: Path,
    *,
    zoom: float = 2.5,
    crop_box: tuple[float, float, float, float] | None = None,
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(pdf_path)
    try:
        page = doc.load_page(page_no_1_based - 1)
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        tmp = out_path.with_suffix(".tmp.png")
        pix.save(tmp.as_posix())

        if crop_box is None:
            tmp.replace(out_path)
            return out_path

        with Image.open(tmp) as im:
            w, h = im.size
            l, t, r, b = crop_box
            crop = (
                int(max(0, min(w, l * w))),
                int(max(0, min(h, t * h))),
                int(max(0, min(w, r * w))),
                int(max(0, min(h, b * h))),
            )
            cropped = im.crop(crop)
            cropped.save(out_path)
        tmp.unlink(missing_ok=True)
        return out_path
    finally:
        doc.close()


def set_run_font(run, *, size_pt: int, bold: bool = False, color: RGBColor | None = None) -> None:
    run.font.name = THEME.font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_title(
    slide,
    title: str,
    *,
    left=Inches(0.8),
    top=Inches(0.4),
    width=Inches(12.0),
    height=Inches(0.9),
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run_font(run, size_pt=28, bold=True, color=THEME.text)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        top + height + Inches(0.05),
        Inches(12.0),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME.accent
    line.line.fill.background()


def add_bullets(
    slide,
    bullets: list[str],
    *,
    left=Inches(0.9),
    top=Inches(1.6),
    width=Inches(6.3),
    height=Inches(5.5),
    font_size=18,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        p.space_before = Pt(0)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            set_run_font(run, size_pt=font_size, color=THEME.text)


def add_kv(
    slide,
    items: list[tuple[str, str]],
    *,
    left=Inches(0.9),
    top=Inches(1.6),
    width=Inches(6.3),
    row_height=Inches(0.75),
) -> None:
    for idx, (k, v) in enumerate(items):
        y = top + row_height * idx
        box = slide.shapes.add_textbox(left, y, width, row_height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        rk = p.add_run()
        rk.text = f"{k}  "
        set_run_font(rk, size_pt=16, bold=True, color=THEME.accent_dark)

        rv = p.add_run()
        rv.text = v
        set_run_font(rv, size_pt=16, color=THEME.text)


def add_picture_fit(
    slide,
    image_path: Path,
    *,
    left,
    top,
    width,
    height,
    add_border: bool = True,
) -> None:
    pic = slide.shapes.add_picture(image_path.as_posix(), left, top, width=width, height=height)
    if add_border:
        pic.line.color.rgb = THEME.line
        pic.line.width = Pt(1)


def build_ppt(
    output_path: Path,
    *,
    project_pdf: Path,
    server_policy_pdf: Path,
    resume_pdf: Path | None = None,
    name: str = "",
    email: str = "",
    github: str = "",
) -> None:
    assets_dir = ROOT / "portfolio_assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    if not project_pdf.exists():
        raise FileNotFoundError(f"project pdf not found: {project_pdf}")
    if not server_policy_pdf.exists():
        raise FileNotFoundError(f"server policy pdf not found: {server_policy_pdf}")

    pdf_project = project_pdf
    pdf_server_policy = server_policy_pdf
    pdf_resume = resume_pdf if (resume_pdf and resume_pdf.exists()) else None

    cover_img = render_pdf_page(pdf_project, 1, assets_dir / "project_cover.png", zoom=2.5)
    network_logical = render_pdf_page(pdf_project, 10, assets_dir / "network_logical.png", zoom=2.5)
    network_physical = render_pdf_page(pdf_project, 11, assets_dir / "network_physical.png", zoom=2.5)
    server_arch = render_pdf_page(pdf_project, 23, assets_dir / "server_arch.png", zoom=2.5)
    packet_table = render_pdf_page(
        pdf_project, 39, assets_dir / "packet_filters.png", zoom=2.8, crop_box=(0.03, 0.38, 0.97, 0.98)
    )
    ssh_policy = render_pdf_page(
        pdf_server_policy, 3, assets_dir / "ssh_policy.png", zoom=2.7, crop_box=(0.04, 0.05, 0.96, 0.98)
    )
    resume_project = None
    if pdf_resume:
        resume_project = render_pdf_page(
            pdf_resume, 3, assets_dir / "resume_project.png", zoom=2.6, crop_box=(0.03, 0.08, 0.97, 0.98)
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Slide 1: Cover
    slide = prs.slides.add_slide(blank_layout)
    add_picture_fit(slide, cover_img, left=Inches(0), top=Inches(0), width=prs.slide_width, height=prs.slide_height, add_border=False)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = rgb("#000000")
    overlay.fill.transparency = 0.55
    overlay.line.fill.background()

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(11.6), Inches(3.6))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME.bg
    box.fill.transparency = 0.08
    box.line.color.rgb = THEME.accent
    box.line.width = Pt(2)

    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(1.3))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "정보보안 · 네트워크 엔지니어\n포트폴리오"
    set_run_font(r, size_pt=44, bold=True, color=THEME.text)

    sub_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(11.0), Inches(1.4))
    tf = sub_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = name or DEFAULT_NAME
    set_run_font(r1, size_pt=26, bold=True, color=THEME.accent_dark)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "기관 인프라 구축 · 모의해킹 기반 취약점 진단 · 탐지/대응 정책 수립"
    set_run_font(r2, size_pt=16, color=THEME.muted)

    # Slide 2: Summary
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "한눈에 보기")
    add_bullets(
        slide,
        [
            "목표: 네트워크 엔지니어(보안/관제 역량 포함) — 구축·운영·탐지·대응까지 End-to-End 경험",
            "핵심: 패킷/로그 기반 이상징후 식별 → 정책(ACL/iptables/UFW/IDS Rule)으로 대응 반영",
            "대표 프로젝트: 기관 인프라 구축 및 모의해킹 기반 통합 보안 취약점 진단/조치",
            "강점: 기획·문서화·협업 리딩(전 직장 팀장 경험) + 기술 구현(네트워크/서버/보안)",
        ],
        width=Inches(12.0),
        height=Inches(5.4),
        font_size=20,
    )

    # Slide 3: Profile
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "프로필")
    add_kv(
        slide,
        [
            ("지원 분야", "네트워크 엔지니어 / 정보보안(관제·대응)"),
            ("교육", "더조은컴퓨터아카데미 강남 — 정보보안 취업캠프 (2025.08~2026.03)"),
            ("학력", "우석대학교 문예창작학과 졸업 (2011.03~2020.08)"),
            ("자격", "리눅스마스터 2급(2026-01-02), 네트워크관리사 2급(2025-11-25)"),
            ("경력", "인포커스 기획팀 팀장 (2023.03~2025.07, 28개월) — 콘텐츠/사업/MICE 기획·운영"),
        ],
        width=Inches(11.8),
    )

    # Slide 4: Skills
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "기술 역량")
    left = Inches(0.9)
    top = Inches(1.5)
    col_w = Inches(6.0)
    row_h = Inches(2.0)

    def add_skill_card(x, y, title, bullets):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("#F5F3FF")
        card.line.color.rgb = THEME.line
        card.line.width = Pt(1)

        t = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), col_w - Inches(0.6), Inches(0.5))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run_font(r, size_pt=18, bold=True, color=THEME.accent_dark)

        b = slide.shapes.add_textbox(x + Inches(0.35), y + Inches(0.7), col_w - Inches(0.7), row_h - Inches(0.85))
        tf = b.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(2)
            for run in p.runs:
                set_run_font(run, size_pt=14, color=THEME.text)

    add_skill_card(
        left,
        top,
        "네트워크",
        ["EIGRP/RIP/Static", "HSRP/VRRP, Gateway 이중화", "ACL/Distribute-List", "SPAN/Port-Security", "NAT(PAT)"],
    )
    add_skill_card(
        left + col_w + Inches(0.6),
        top,
        "보안·관제",
        ["Wireshark 패킷 분석", "Security Onion + Snort/Sguil", "모의해킹 시나리오 기반 탐지", "iptables/UFW 정책", "취약점 진단(Web/DB/SSH)"],
    )
    add_skill_card(
        left,
        top + row_h + Inches(0.5),
        "서버·서비스",
        ["Linux(Rocky/Ubuntu/CentOS)", "Nginx Reverse Proxy + Apache", "BIND9(DNS), MariaDB", "Postfix/Dovecot/Roundcube", "NFS 기반 백업/로그 수집"],
    )
    add_skill_card(
        left + col_w + Inches(0.6),
        top + row_h + Inches(0.5),
        "자동화·도구",
        ["Python(패킷/로그 분석)", "scapy/requests 활용", "GNS3/Packet Tracer", "MRTG/Cacti 시각화", "문서화/표준 산출물"],
    )

    # Slide 5: Project overview (resume page image + bullets)
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "대표 프로젝트 요약")
    if resume_project:
        add_picture_fit(
            slide,
            resume_project,
            left=Inches(7.1),
            top=Inches(1.55),
            width=Inches(5.35),
            height=Inches(5.55),
            add_border=True,
        )
        bullets_w = Inches(6.1)
    else:
        bullets_w = Inches(12.0)

    add_bullets(
        slide,
        [
            "프로젝트: 기관 인프라 구축 및 모의해킹 기반 통합 보안 취약점 진단/조치",
            "기간: 2025.12.29 ~ 2026.01.13 (약 2주)",
            "핵심 흐름: 설계 → 공격 시나리오 구현(Red) → 탐지·대응(Blue) → 정책 환류",
            "담당: 패킷 식별/분석 및 대응 정책 수립(보안·관제), 산출물 표준화/기획 리딩",
        ],
        width=bullets_w,
        height=Inches(5.6),
        font_size=18,
    )

    # Slide 6: Project goals
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "프로젝트 목표 (기획서 기반)")
    add_bullets(
        slide,
        [
            "네트워크 구축: 회선·Gateway 이중화 + EIGRP 적용, 정책 기반 라우팅 및 접근 제어",
            "관제: MRTG/Cacti로 가시성 확보, Wireshark·IDS로 비정상 트래픽 탐지",
            "보안: 프로토콜 인증/ACL/Distribute-List, Port-security, NAT(PAT)로 내부 보호",
            "대응: Python 기반 패킷 분석 → 룰/ACL 반영, iptables/UFW 기반 차단 정책 적용",
            "서버: 서비스 분산 구축 + 백업/로그 중앙화로 가용성·무결성 확보",
        ],
        width=Inches(12.0),
        height=Inches(5.7),
        font_size=19,
    )

    # Slide 7: Network architecture
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "네트워크 아키텍처")
    add_picture_fit(
        slide,
        network_logical,
        left=Inches(0.9),
        top=Inches(1.55),
        width=Inches(5.95),
        height=Inches(5.55),
        add_border=True,
    )
    add_picture_fit(
        slide,
        network_physical,
        left=Inches(6.95),
        top=Inches(1.55),
        width=Inches(5.95),
        height=Inches(5.55),
        add_border=True,
    )

    # Slide 8: Server architecture
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "서버 구성 및 서비스 분산")
    add_picture_fit(
        slide,
        server_arch,
        left=Inches(6.9),
        top=Inches(1.55),
        width=Inches(6.05),
        height=Inches(5.55),
        add_border=True,
    )
    add_bullets(
        slide,
        [
            "역할 분리: Web(Reverse Proxy) / WAS(CMS) / DBMS / DNS / Mail / WHD / Monitoring / Backup",
            "모니터링: Cacti(CPU/MEM/Disk), MRTG(대역폭/상태)로 조기 징후 탐지",
            "백업/로그: NFS 기반 중앙 저장 + 주기 수집(Crontab)으로 무결성·가용성 강화",
        ],
        width=Inches(6.0),
        height=Inches(5.5),
        font_size=18,
    )

    # Slide 9: Vulnerability -> Policy
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "취약점 → 보안정책 (예시)")

    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.55), Inches(5.95), Inches(5.55))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = rgb("#FEF2F2")
    left_card.line.color.rgb = THEME.line
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.55), Inches(5.95), Inches(5.55))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = rgb("#ECFDF5")
    right_card.line.color.rgb = THEME.line

    t1 = slide.shapes.add_textbox(Inches(1.2), Inches(1.75), Inches(5.5), Inches(0.5))
    tf = t1.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "의도된 취약 설정(진단 목적)"
    set_run_font(r, size_pt=18, bold=True, color=rgb("#B91C1C"))

    t2 = slide.shapes.add_textbox(Inches(7.25), Inches(1.75), Inches(5.5), Inches(0.5))
    tf = t2.text_frame
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = "보안 구축/개선 정책"
    set_run_font(r, size_pt=18, bold=True, color=rgb("#047857"))

    add_bullets(
        slide,
        [
            "SNMP rocommunity public (전체 접근)",
            "Cacti HTTP(평문)·세션 탈취 위험",
            "NFS no_root_squash(권한 상승)",
            "Roundcube 브루트포스 제한 없음",
            "MariaDB 계정 허용 IP ‘%’",
        ],
        left=Inches(1.2),
        top=Inches(2.35),
        width=Inches(5.5),
        height=Inches(4.6),
        font_size=16,
    )
    add_bullets(
        slide,
        [
            "SNMP 접근 IP 제한(커뮤니티/ACL)",
            "HTTPS 적용·기본 admin 비활성화",
            "root_squash 적용·관리망 한정 통신",
            "SSH 공개키 인증·root 로그인 차단",
            "iptables/UFW 화이트리스트 정책",
        ],
        left=Inches(7.25),
        top=Inches(2.35),
        width=Inches(5.5),
        height=Inches(4.6),
        font_size=16,
    )

    # Slide 10: Detection & Response (packet table + policy image)
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "패킷 식별 · 탐지 룰 · 차단 가이드")
    add_picture_fit(
        slide,
        packet_table,
        left=Inches(0.9),
        top=Inches(1.55),
        width=Inches(7.0),
        height=Inches(5.55),
        add_border=True,
    )
    add_picture_fit(
        slide,
        ssh_policy,
        left=Inches(8.15),
        top=Inches(1.55),
        width=Inches(4.75),
        height=Inches(5.55),
        add_border=True,
    )

    # Slide 11: Outcomes
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "기대효과")
    add_bullets(
        slide,
        [
            "분산 구축으로 네트워크/서버 안정성과 가용성 확보",
            "로그 분석·백업으로 데이터 손실 위협 대응 및 빠른 복구",
            "통합 모니터링으로 선제적 위협 탐지 및 즉각 대응",
            "화이트리스트 기반 접근 제어 + 구간 암호화(SSL)로 데이터 유출 최소화",
            "Reverse Proxy 구성으로 부하 분산 및 사용자 응답성 개선",
        ],
        width=Inches(12.0),
        height=Inches(5.6),
        font_size=20,
    )

    # Slide 12: Deliverables
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "산출물")
    add_bullets(
        slide,
        [
            "2차 프로젝트 기획안 v0.4.1(퍼플).pdf — 초기 요구사항/구성도/테스트·검증 방향",
            "2차 프로젝트 기획안 v0.9.pdf — 최신 버전(네트워크·서버·모의해킹·로그/패킷 분석·정책) 정리",
            "네트워크 보안관제 보고서.pdf — 시나리오별 Wireshark 필터/식별 기준 정리",
            "서버 보안 정책 문서(미완성).pdf — 로그 제출 기준 및 iptables/UFW 차단 가이드(초안)",
            "이력서(PDF) — 포함 시 Slide 5 자동 삽입(경로: _private/resume/...)",
        ],
        width=Inches(12.0),
        height=Inches(5.6),
        font_size=18,
    )

    # Slide 13: Career highlight
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "업무 역량 (전 직장 경험 기반)")
    add_bullets(
        slide,
        [
            "팀장 경험: 기획총괄·아트디렉션·대외 커뮤니케이션·일정/리스크 관리",
            "사업 성과: 2년간 입찰/지원사업 총 10.8억 규모 성과 주도, 22건(4.55억) 수행",
            "현장 운영: 국제회의/컨벤션 등 대형 행사 기획부터 운영까지 End-to-End 수행",
            "보안/네트워크와의 연결: 문서 표준화·현업 협업·운영 관점의 정책 설계/전달 역량",
        ],
        width=Inches(12.0),
        height=Inches(5.6),
        font_size=20,
    )

    # Slide 14: Contact
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "연락처")
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.1), Inches(10.4), Inches(3.0))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = rgb("#F9FAFB")
    contact_box.line.color.rgb = THEME.line
    contact_box.line.width = Pt(1)

    box = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.5), Inches(2.2))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = name or DEFAULT_NAME
    set_run_font(r, size_pt=28, bold=True, color=THEME.accent_dark)

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"E-mail: {email}" if email else "E-mail: (옵션)"
    set_run_font(r, size_pt=18, color=THEME.text)

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = f"GitHub: {github}" if github else "GitHub: (옵션)"
    set_run_font(r, size_pt=18, color=THEME.text)

    prs.save(output_path.as_posix())


DEFAULT_PROJECT_PDF = REPO_ROOT / "docs" / "project-proposals" / "2차 프로젝트 기획안 v0.9.pdf"
DEFAULT_SERVER_POLICY_PDF = REPO_ROOT / "reports" / "security-policy" / "서버 보안 정책 문서(미완성).pdf"
DEFAULT_RESUME_PDF = REPO_ROOT / "_private" / "resume" / "이력서_정보보안_박신우.pdf"
DEFAULT_OUT = REPO_ROOT / "profile" / "Portfolio.pptx"
DEFAULT_NAME = "박신우"
DEFAULT_GITHUB = "https://github.com/SINUNET1012"


def main() -> None:
    parser = argparse.ArgumentParser(description="PDF 기반 포트폴리오 PPTX 생성기")
    parser.add_argument("--project-pdf", type=Path, default=DEFAULT_PROJECT_PDF, help="프로젝트 기획안 PDF 경로")
    parser.add_argument("--server-policy-pdf", type=Path, default=DEFAULT_SERVER_POLICY_PDF, help="서버 보안 정책 PDF 경로")
    parser.add_argument("--resume-pdf", type=Path, default=DEFAULT_RESUME_PDF, help="(선택) 이력서 PDF 경로")
    parser.add_argument("--out", type=Path, default=DEFAULT_OUT, help="출력 PPTX 경로")
    parser.add_argument("--name", default=DEFAULT_NAME, help="표기 이름(슬라이드 내)")
    parser.add_argument("--email", default="", help="(선택) 이메일")
    parser.add_argument("--github", default=DEFAULT_GITHUB, help="(선택) GitHub 링크")
    args = parser.parse_args()

    name = str(args.name or "").strip()
    email = str(args.email or "").strip()
    github = str(args.github or "").strip()

    if not args.project_pdf.exists():
        raise SystemExit(f"project pdf not found: {args.project_pdf}")
    if not args.server_policy_pdf.exists():
        raise SystemExit(f"server policy pdf not found: {args.server_policy_pdf}")

    out = args.out
    out.parent.mkdir(parents=True, exist_ok=True)
    build_ppt(
        out,
        project_pdf=args.project_pdf,
        server_policy_pdf=args.server_policy_pdf,
        resume_pdf=args.resume_pdf,
        name=name,
        email=email,
        github=github,
    )
    print(out.as_posix())


if __name__ == "__main__":
    main()
